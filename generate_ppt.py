"""Generate a professional project presentation PPT in Chinese."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BLUE_DARK = RGBColor(0x0D, 0x1B, 0x2A)
BLUE_MID = RGBColor(0x1B, 0x2A, 0x4A)
BLUE_ACCENT = RGBColor(0x41, 0x5A, 0x77)
BLUE_LIGHT = RGBColor(0x77, 0x8D, 0xA9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF0, 0xC0, 0x40)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
GREEN = RGBColor(0x00, 0xE6, 0x76)
RED_SOFT = RGBColor(0xFF, 0x6B, 0x6B)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_SUB = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
TEAL = RGBColor(0x14, 0xB8, 0xA6)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0


def add_shape_bg(slide, color, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=GRAY_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=GRAY_TEXT, bullet_color=None, spacing=Pt(6),
                    font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             accent_color=CYAN, bg_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top, Inches(0.06), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    add_text_box(slide, left + Inches(0.2), top + Inches(0.08),
                 width - Inches(0.3), Inches(0.35),
                 title, font_size=13, color=accent_color, bold=True)

    y = top + Inches(0.42)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.2), y,
                     width - Inches(0.3), Inches(0.25),
                     line, font_size=10, color=GRAY_SUB)
        y += Inches(0.22)


def add_page_number(slide, num, total):
    add_text_box(slide, Inches(8.8), Inches(6.8), Inches(1), Inches(0.3),
                 f"{num}/{total}", font_size=9, color=GRAY_LIGHT,
                 alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, number, title, subtitle=""):
    add_gradient_bg(slide, BLUE_DARK, BLUE_MID)

    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(8), Inches(0.6),
                 f"PART {number}", font_size=16, color=BLUE_LIGHT,
                 bold=True, font_name="Consolas")

    add_text_box(slide, Inches(0.8), Inches(2.6), Inches(8), Inches(1.0),
                 title, font_size=36, color=WHITE, bold=True)

    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(3.6), Inches(8), Inches(0.6),
                     subtitle, font_size=16, color=BLUE_LIGHT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(3.45), Inches(1.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    TOTAL = 20

    # ========================================================
    # SLIDE 1: Cover
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, BLUE_DARK, RGBColor(0x15, 0x25, 0x3B))

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0), Inches(0), Inches(10), Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = CYAN
    top_bar.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(0.5),
                 "AI PHYSICS EXPERIMENT PLATFORM", font_size=14,
                 color=CYAN, bold=True, font_name="Consolas")

    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(8.5), Inches(1.2),
                 "AI 物理实验仿真平台", font_size=42,
                 color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(8.5), Inches(0.8),
                 "基于 NVIDIA Isaac Sim + PhysX 5 的\nGPU 加速全栈物理实验系统",
                 font_size=18, color=BLUE_LIGHT)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(4.6), Inches(2.0), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = GOLD
    sep.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(4.9), Inches(8.5), Inches(0.4),
                 "项目汇报  |  2026 年春季学期", font_size=14, color=BLUE_LIGHT)

    tech_items = ["React 19 + TypeScript", "WebRTC 实时视频", "WebSocket 双向通信",
                  "PhysX 5 物理引擎", "NVIDIA RTX 5090"]
    x = Inches(0.8)
    for item in tech_items:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, Inches(5.8), Inches(1.6), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE_ACCENT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(9)
        p.font.color.rgb = CYAN
        p.font.name = "Consolas"
        p.alignment = PP_ALIGN.CENTER
        x += Inches(1.7)

    add_page_number(slide, 1, TOTAL)

    # ========================================================
    # SLIDE 2: Table of Contents
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5),
                 "目录 CONTENTS", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.1), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    sections = [
        ("01", "项目概览", "平台定位、系统架构与技术栈"),
        ("02", "我的核心工作", "实验一全栈开发 + 整体框架搭建"),
        ("03", "后端开发", "WebRTC 服务器、视频管线与遥测系统"),
        ("04", "前端开发", "React 组件体系与前端适配挑战"),
        ("05", "视频流投送", "如何把 Isaac Sim 画面推到浏览器"),
        ("06", "同学代码集成", "实验二（大摆角摆）与实验七（动量守恒）"),
        ("07", "技术挑战与解决", "踩过的坑与工程解决方案"),
        ("08", "项目成果与展望", "成果展示与未来规划"),
    ]

    y = Inches(1.6)
    for num, title, desc in sections:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.8), y, Inches(0.7), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE_DARK
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.color.rgb = CYAN
        p.font.bold = True
        p.font.name = "Consolas"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(6)

        add_text_box(slide, Inches(1.7), y + Inches(0.02), Inches(3.5), Inches(0.3),
                     title, font_size=16, color=BLUE_DARK, bold=True)
        add_text_box(slide, Inches(1.7), y + Inches(0.30), Inches(6), Inches(0.25),
                     desc, font_size=11, color=GRAY_SUB)
        y += Inches(0.7)

    add_page_number(slide, 2, TOTAL)

    # ========================================================
    # SLIDE 3: Section — 项目概览
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "01", "项目概览", "平台定位 · 系统架构 · 技术栈")
    add_page_number(slide, 3, TOTAL)

    # ========================================================
    # SLIDE 4: 项目背景与目标
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "项目背景与目标", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8.5), Inches(0.8),
                 "构建一个基于 GPU 加速的全栈物理实验仿真平台，让学生通过浏览器即可\n"
                 "操作大学物理实验，获得实时 3D 仿真画面、遥测数据和自动实验报告。",
                 font_size=14, color=GRAY_TEXT)

    cards_data = [
        ("硬件平台", ["NVIDIA RTX 5090 GPU 服务器", "Ubuntu Linux 大学机房部署", "远程 SSH + 浏览器访问"], CYAN),
        ("仿真引擎", ["NVIDIA Isaac Sim (Kit 应用)", "PhysX 5 GPU 物理引擎", "USD 场景描述格式"], PURPLE),
        ("Web 前端", ["React 19 + TypeScript", "Vite 构建 + Tailwind CSS", "Recharts 实时数据可视化"], GREEN),
        ("通信协议", ["WebRTC 实时视频流", "WebSocket 双向控制/遥测", "HTTP REST 信令接口"], ORANGE),
    ]

    x = Inches(0.6)
    for title, lines, color in cards_data:
        add_card(slide, x, Inches(2.5), Inches(2.1), Inches(1.8), title, lines, color)
        x += Inches(2.3)

    stats = [("8", "个物理实验"), ("3", "个已完成"), ("2136", "行后端代码"),
             ("1237", "行前端代码"), ("3", "个通信端口")]
    x = Inches(0.6)
    for val, label in stats:
        add_text_box(slide, x, Inches(4.7), Inches(1.5), Inches(0.5),
                     val, font_size=28, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(5.2), Inches(1.5), Inches(0.3),
                     label, font_size=11, color=GRAY_SUB, alignment=PP_ALIGN.CENTER)
        x += Inches(1.8)

    add_page_number(slide, 4, TOTAL)

    # ========================================================
    # SLIDE 5: 系统架构图
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "系统架构", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    layers = [
        ("浏览器层 (React 19 + TypeScript + Vite + Tailwind)",
         "Landing.tsx  |  LevelSelect.tsx  |  ExperimentView.tsx  |  WebRTCIsaacViewer.tsx\n"
         "experiments.ts (8个实验UI定义)  |  isaacService.ts (WebSocket客户端)  |  config.ts (自动IP)",
         RGBColor(0x22, 0x7C, 0x9D), Inches(1.3)),
        ("通信层 (WebSocket :30000 + WebRTC :8080 + HTTP REST)",
         "WebSocket: 实验控制命令 + 100Hz遥测广播  |  WebRTC: 实时视频帧传输 (30fps)\n"
         "HTTP POST: SDP信令交换 + 相机控制  |  WS-JPEG: SSH隧道降级方案",
         RGBColor(0xE9, 0x72, 0x2D), Inches(3.0)),
        ("Isaac Sim Python 运行时 (core/webrtc_server.py)",
         "WebRTCServer: aiohttp + aiortc 异步服务  |  IsaacSimVideoTrack: 视口帧捕获\n"
         "CameraController: 轨道/平移/缩放  |  实验状态机 + 碰撞检测 + RK4积分",
         RGBColor(0x26, 0x4E, 0x53), Inches(4.7)),
        ("物理引擎层 (PhysX 5 + USD Scene)",
         "Experiment/exp.usd: 统一场景  |  DynamicCuboid / FixedCuboid / VisualCuboid\n"
         "PhysicsMaterial: 摩擦/恢复系数  |  RigidBodyAPI + MassAPI + RevoluteJoint",
         RGBColor(0x6B, 0x21, 0xA8), Inches(6.1)),
    ]

    for title, desc, color, y in layers:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), y, Inches(9), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        p2.font.name = "Consolas"
        p2.space_before = Pt(6)

    for y in [Inches(2.85), Inches(4.55), Inches(5.95) + Pt(3)]:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                      Inches(4.8), y, Inches(0.4), Inches(0.25))
        arr.fill.solid()
        arr.fill.fore_color.rgb = GRAY_LIGHT
        arr.line.fill.background()

    add_page_number(slide, 5, TOTAL)

    # ========================================================
    # SLIDE 6: Section — 我的核心工作
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "02", "我的核心工作", "实验一全栈开发 · 整体框架搭建 · 代码集成")
    add_page_number(slide, 6, TOTAL)

    # ========================================================
    # SLIDE 7: 工作量总览 (rewritten: accessible, non-jargon)
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "个人工作总览", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    # --- Three vertical cards side by side ---
    card_data = [
        {
            "tag": "01", "tag_color": CYAN,
            "title": "实验一 · 角动量守恒",
            "subtitle": "独立完成全部代码",
            "items": [
                "搭建完整的 3D 虚拟实验场景，还原\n真实转台、环和盘的实验器材外观",
                "编写物理仿真逻辑：转盘旋转、物体\n下落、守恒过程的全流程模拟",
                "服务器实时采集角速度、角动量、动能\n等物理量，持续推送到浏览器",
                "设计前端【设置-旋转-落下-记录】\n4 步引导式流程，共 4 个试次",
                "实时曲线图 + 数据表 + CSV 导出\n+ 一键生成完整 PDF 实验报告",
                "另有命令行离线模式，可批量生成\n数据、图表和 Markdown 分析报告",
            ],
            "number": "694+", "number_label": "行仿真代码",
        },
        {
            "tag": "02", "tag_color": PURPLE,
            "title": "全栈平台框架",
            "subtitle": "从零搭建整个系统骨架",
            "items": [
                "独立编写 2100+ 行后端服务器代码，\n连接仿真引擎与浏览器端",
                "实现 GPU 渲染画面的实时视频推流，\n让用户在浏览器中看到 3D 仿真",
                "设计备用视频通道：网络受限时自动\n无缝切换，保证画面不中断",
                "搭建浏览器与服务器的双向实时通信\n（发送控制指令 + 回传物理数据）",
                "基于学长旧前端代码做大幅重构改造，\n解决 10+ 个兼容性和适配问题",
                "设计统一配置系统，让 8 个实验\n共享一套前端框架和交互逻辑",
            ],
            "number": "2100+", "number_label": "行后端代码",
        },
        {
            "tag": "03", "tag_color": GREEN,
            "title": "同学代码集成",
            "subtitle": "接入实验二和实验七",
            "items": [
                "将两位同学的独立脚本改造为可在\n网页端实时操控的在线实验",
                "实验二（大摆角摆）：重建 3D 场景，\n集成摆动计算，实时显示角度和周期",
                "实验七（碰撞实验）：重建小车碰撞\n场景，检测碰撞并精确捕获速度变化",
                "两个实验都需要重写场景搭建方式\n来兼容实时视频推流管线",
                "两个实验都配套设计了完整的\n4 试次实验流程 + PDF 报告功能",
                "涉及大量联调与适配工作，确保\n物理精度和显示效果的一致性",
            ],
            "number": "2", "number_label": "个实验接入",
        },
    ]

    card_w = Inches(2.9)
    card_gap = Inches(0.2)
    card_x_start = Inches(0.4)
    card_top = Inches(1.25)

    for ci, cd in enumerate(card_data):
        cx = card_x_start + ci * (card_w + card_gap)

        # Card background
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_top,
            card_w, Inches(5.65))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = WHITE
        card_bg.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card_bg.line.width = Pt(1)

        # Top accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx, card_top,
            card_w, Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = cd["tag_color"]
        accent_bar.line.fill.background()

        # Tag number circle
        tag_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            cx + Inches(0.15), card_top + Inches(0.22),
            Inches(0.42), Inches(0.42))
        tag_circle.fill.solid()
        tag_circle.fill.fore_color.rgb = cd["tag_color"]
        tag_circle.line.fill.background()
        tf = tag_circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = cd["tag"]
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Consolas"
        p.alignment = PP_ALIGN.CENTER

        # Title
        add_text_box(slide, cx + Inches(0.65), card_top + Inches(0.2),
                     card_w - Inches(0.8), Inches(0.3),
                     cd["title"], font_size=14, color=BLUE_DARK, bold=True)
        # Subtitle
        add_text_box(slide, cx + Inches(0.65), card_top + Inches(0.48),
                     card_w - Inches(0.8), Inches(0.25),
                     cd["subtitle"], font_size=10, color=cd["tag_color"], bold=True)

        # Divider
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            cx + Inches(0.15), card_top + Inches(0.80),
            card_w - Inches(0.3), Pt(1.5))
        div.fill.solid()
        div.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        div.line.fill.background()

        # Bullet items
        item_y = card_top + Inches(0.92)
        for item_text in cd["items"]:
            # Bullet dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                cx + Inches(0.18), item_y + Inches(0.06),
                Inches(0.07), Inches(0.07))
            dot.fill.solid()
            dot.fill.fore_color.rgb = cd["tag_color"]
            dot.line.fill.background()

            add_text_box(slide, cx + Inches(0.32), item_y,
                         card_w - Inches(0.45), Inches(0.55),
                         item_text, font_size=9, color=GRAY_SUB)
            item_y += Inches(0.58)

        # Bottom number highlight
        num_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cx + Inches(0.3), card_top + Inches(4.95),
            card_w - Inches(0.6), Inches(0.55))
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = cd["tag_color"]
        num_bg.line.fill.background()

        num_tf = num_bg.text_frame
        num_tf.word_wrap = True
        p_num = num_tf.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        run_val = p_num.add_run()
        run_val.text = cd["number"]
        run_val.font.size = Pt(20)
        run_val.font.color.rgb = WHITE
        run_val.font.bold = True
        run_val.font.name = "Consolas"
        run_label = p_num.add_run()
        run_label.text = "  " + cd["number_label"]
        run_label.font.size = Pt(10)
        run_label.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
        run_label.font.bold = False
        run_label.font.name = "Microsoft YaHei"

    add_page_number(slide, 7, TOTAL)

    # ========================================================
    # SLIDE 7b: 工作量数据一览 (metrics page)
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.5),
                 "工作量数据一览", font_size=28, color=BLUE_DARK, bold=True)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(0.92), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    # ---- Hero numbers row ----
    hero_nums = [
        ("10,000+",  "行代码总量",    CYAN),
        ("5,243",    "行前端代码",    RGBColor(0x3B, 0x82, 0xF6)),
        ("4,800+",   "行后端代码",    PURPLE),
        ("3",        "个完整实验",    GREEN),
        ("8",        "个实验框架",    ORANGE),
        ("2",        "条视频通道",    TEAL),
    ]
    hero_w = Inches(1.38)
    hero_gap = Inches(0.15)
    hero_x = Inches(0.35)
    hero_y = Inches(1.15)
    for val, label, color in hero_nums:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      hero_x, hero_y, hero_w, Inches(1.05))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        box.line.width = Pt(1)

        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           hero_x, hero_y, hero_w, Inches(0.05))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = color
        top_line.line.fill.background()

        add_text_box(slide, hero_x, hero_y + Inches(0.15),
                     hero_w, Inches(0.45),
                     val, font_size=22, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")
        add_text_box(slide, hero_x, hero_y + Inches(0.62),
                     hero_w, Inches(0.3),
                     label, font_size=10, color=GRAY_SUB,
                     alignment=PP_ALIGN.CENTER)
        hero_x += hero_w + hero_gap

    # ---- "Built" column (left) ----
    col_left_x = Inches(0.5)
    col_right_x = Inches(5.1)
    section_y = Inches(2.45)

    built_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           col_left_x, section_y,
                                           Inches(4.3), Inches(0.38))
    built_header.fill.solid()
    built_header.fill.fore_color.rgb = CYAN
    built_header.line.fill.background()
    tf = built_header.text_frame
    p = tf.paragraphs[0]
    p.text = "  搭建的功能模块"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    built_items = [
        "实验一角动量守恒 — 从建模到报告的完整实验流程",
        "2100+ 行后端服务器 — 视频推流 + 实时通信 + 数据广播",
        "1236 行实验控制界面 — 三种实验各有专属 UI 布局",
        "实时视频推流系统 — 将 GPU 画面送到浏览器播放",
        "备用图片流通道 — 网络受限时自动无缝切换",
        "100Hz 物理数据广播 — 角速度/动量/能量实时推送",
        "4 试次实验流程引擎 — 实验一和实验七各一套",
        "2 套 PDF 报告生成器 — 含图表、数据表、误差分析",
        "CSV 数据导出功能 — 支持离线数据进一步分析",
        "3D 相机交互系统 — 鼠标旋转、平移、缩放视角",
        "8 个实验的统一配置框架 — 控件/图表/参数声明式管理",
        "8 套摄像机预设脚本 — 每个实验独立视角",
    ]

    by = section_y + Inches(0.45)
    for i, item in enumerate(built_items):
        bg_c = RGBColor(0xF0, 0xFD, 0xFA) if i % 2 == 0 else WHITE
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      col_left_x, by, Inches(4.3), Inches(0.3))
        row.fill.solid()
        row.fill.fore_color.rgb = bg_c
        row.line.color.rgb = RGBColor(0xE8, 0xF0, 0xF0)
        row.line.width = Pt(0.5)

        add_text_box(slide, col_left_x + Inches(0.08), by,
                     Inches(4.1), Inches(0.3),
                     "✓  " + item, font_size=9, color=GRAY_TEXT)
        by += Inches(0.3)

    # ---- "Fixed / Solved" column (right) ----
    fixed_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           col_right_x, section_y,
                                           Inches(4.4), Inches(0.38))
    fixed_header.fill.solid()
    fixed_header.fill.fore_color.rgb = ORANGE
    fixed_header.line.fill.background()
    tf = fixed_header.text_frame
    p = tf.paragraphs[0]
    p.text = "  解决的问题与适配工作"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"

    fixed_items = [
        "重构学长前端旧代码 — 大幅改造使其适配新后端",
        "修复硬编码 IP 地址 — 改为自动检测服务器地址",
        "修复前端样式构建方式 — 从 CDN 迁移到工程化构建",
        "修复命令映射缺失 — 补全暂停/重置等缺失的控制逻辑",
        "新增实验切换协议 — 原系统无法在多个实验间切换",
        "解决视频连接不稳定 — 设计自动降级备用方案",
        "解决摄像机被引擎重置 — 多次延迟重设对抗异步初始化",
        "解决帧捕获与仿真冲突 — 重写渲染方式使两者兼容",
        "解决物理引擎启动吞速度 — 添加暖机等待阶段",
        "解决碰撞精度不足 — 调高求解器精度参数",
        "解决帧捕获首次返回空 — 预热机制 + 自动重建",
        "将两位同学脚本改造接入 — 重写场景 + 适配管线 + UI定制",
    ]

    fy = section_y + Inches(0.45)
    for i, item in enumerate(fixed_items):
        bg_c = RGBColor(0xFF, 0xF7, 0xED) if i % 2 == 0 else WHITE
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      col_right_x, fy, Inches(4.4), Inches(0.3))
        row.fill.solid()
        row.fill.fore_color.rgb = bg_c
        row.line.color.rgb = RGBColor(0xF0, 0xE8, 0xE0)
        row.line.width = Pt(0.5)

        add_text_box(slide, col_right_x + Inches(0.08), fy,
                     Inches(4.2), Inches(0.3),
                     "✦  " + item, font_size=9, color=GRAY_TEXT)
        fy += Inches(0.3)

    # ---- Bottom summary bar ----
    summary_y = Inches(6.25)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.5), summary_y,
                                  Inches(9.0), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_DARK
    bar.line.fill.background()

    summary_texts = [
        ("12", "项功能搭建"),
        ("12", "个问题解决"),
        ("6", "个前端页面组件"),
        ("2", "份 PDF 报告模板"),
        ("3", "种通信协议"),
    ]
    sx = Inches(0.7)
    for val, label in summary_texts:
        add_text_box(slide, sx, summary_y + Inches(0.02),
                     Inches(0.5), Inches(0.45),
                     val, font_size=18, color=GOLD, bold=True,
                     alignment=PP_ALIGN.RIGHT, font_name="Consolas")
        add_text_box(slide, sx + Inches(0.55), summary_y + Inches(0.08),
                     Inches(1.3), Inches(0.35),
                     label, font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC),
                     alignment=PP_ALIGN.LEFT)
        sx += Inches(1.75)

    add_page_number(slide, 8, TOTAL)

    # ========================================================
    # SLIDE 8: 实验一 — 物理模型
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "实验一 · 角动量守恒 — 物理建模", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8.5), Inches(0.6),
                 "将环（或盘）落到旋转转台上，验证角动量守恒定律 L = Iω = const",
                 font_size=14, color=GRAY_TEXT)

    physics_cards = [
        ("物理原理", [
            "I_initial × ω_i = I_final × ω_f",
            "I_disk = ½mR²  (实心圆盘)",
            "I_ring = ½m(R₁² + R₂²) + mx²",
            "平行轴定理: I = I_cm + md²",
            "摩擦能量损失: 1-3%",
        ], CYAN),
        ("仿真实现", [
            "DynamicCuboid 代理刚体驱动物理",
            "RevoluteJoint 约束Z轴旋转",
            "MassAPI 覆写精确惯量张量",
            "PhysicsMaterial 接触摩擦",
            "随机偏心 σ≈2mm 模拟人手误差",
        ], PURPLE),
        ("数据采集", [
            "DynamicControl API 读角速度",
            "100Hz 遥测广播 (角速度/角动量/KE)",
            "碰撞前后角动量百分比差",
            "动能损失百分比计算",
            "传感器噪声 ±0.002 rad/s",
        ], GREEN),
        ("USD 场景", [
            "加载队友3D模型 exp1.usd",
            "自动发现 disk/ring prim",
            "剥离嵌入物理(防冲突)",
            "代理体 ↔ 视觉体 逐帧同步",
            "旋转标记点可视化",
        ], ORANGE),
    ]

    x = Inches(0.4)
    for title, items, color in physics_cards:
        add_card(slide, x, Inches(2.1), Inches(2.2), Inches(2.5), title, items, color)
        x += Inches(2.35)

    add_text_box(slide, Inches(0.8), Inches(4.9), Inches(8.5), Inches(0.5),
                 "核心公式:  Iω_before = Iω_after  →  ω_f = I_i/(I_i + I_drop) × ω_i × (1-μ_friction)",
                 font_size=13, color=BLUE_DARK, bold=True, font_name="Consolas")

    formulas = [
        "I_RI = ½ × m_disk × R² + ½ × m_pulley × R_p²",
        "F_RI = I_RI + ½ × m_ring × (R₁² + R₂²) + m_ring × x²",
        "动能损失: ΔKE/KE_i = 1 - (I_i/I_f) → 理论 ~30-50%",
    ]
    y = Inches(5.5)
    for f in formulas:
        add_text_box(slide, Inches(1.0), y, Inches(8), Inches(0.25),
                     f, font_size=10, color=GRAY_SUB, font_name="Consolas")
        y += Inches(0.25)

    add_page_number(slide, 8, TOTAL)

    # ========================================================
    # SLIDE 9: 实验一 — 前端4试次状态机
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "实验一 · 前端4试次状态机", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    states = [
        ("IDLE\n设置参数", RGBColor(0x64, 0x74, 0x8B)),
        ("SPINNING\n转盘旋转", RGBColor(0x10, 0xB9, 0x81)),
        ("DROPPED\n落下物体", RGBColor(0xF5, 0x9E, 0x0B)),
        ("RECORDED\n记录数据", RGBColor(0x3B, 0x82, 0xF6)),
    ]

    x = Inches(0.5)
    for i, (label, color) in enumerate(states):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, Inches(1.4), Inches(1.8), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if i < 3:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          x + Inches(1.85), Inches(1.7),
                                          Inches(0.4), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = GRAY_LIGHT
            arr.line.fill.background()
        x += Inches(2.3)

    features = [
        ("左侧面板", [
            "物理常数表 (Ring/Lower Disk/Upper Disk/Pulley 精确参数)",
            "初始角速度滑块 ω₀ = 15~30 rad/s",
            "落下物体选择器 (Ring / Upper Disk)",
            "质量滑块 (首次 Spin 后锁定，保证实验一致性)",
            "当前试次信息卡片 (FRI/FAV/偏移量)",
            "实时角速度折线图 (Recharts LineChart)",
        ], CYAN),
        ("底部面板", [
            "步骤指示器: idle → spinning (3s倒计时) → dropped → recorded",
            "4步操作按钮: Spin → Drop → Record → Next Trial",
            "试次数据表 (13列: IAV/FAV/x/IRI/FRI/IAM/FAM/%diff/InitK/FinalK/Energy%)",
            "CSV 导出 + PDF 实验报告一键生成",
        ], PURPLE),
        ("PDF 报告", [
            "使用 @react-pdf/renderer 纯前端生成",
            "Canvas 渲染每个试次的角速度变化图",
            "包含: 物理常数、计算公式、4试次数据表、误差分析",
            "renderTrialChart() 自定义图表绘制器",
        ], GREEN),
    ]

    y = Inches(2.6)
    for title, items, color in features:
        add_text_box(slide, Inches(0.8), y, Inches(3), Inches(0.3),
                     "◆ " + title, font_size=13, color=color, bold=True)
        y += Inches(0.32)
        for item in items:
            add_text_box(slide, Inches(1.2), y, Inches(8), Inches(0.23),
                         "· " + item, font_size=10, color=GRAY_SUB)
            y += Inches(0.24)
        y += Inches(0.1)

    add_page_number(slide, 9, TOTAL)

    # ========================================================
    # SLIDE 10: Section — 后端开发
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "03", "后端开发", "WebRTC 服务器 · 视频管线 · 遥测系统")
    add_page_number(slide, 10, TOTAL)

    # ========================================================
    # SLIDE 11: WebRTC 服务器
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "后端核心 — core/webrtc_server.py", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.5), Inches(0.4),
                 "2100+ 行 Python 异步服务器，运行在 Isaac Sim 进程内部",
                 font_size=14, color=GRAY_TEXT)

    backend_cards = [
        ("WebRTCServer 主类", [
            "aiohttp 异步 HTTP + WebSocket",
            "aiortc WebRTC 信令处理",
            "8个实验的状态变量管理",
            "SDP offer/answer + ICE patch",
            "CORS 跨域 + 静态文件服务",
        ], CYAN),
        ("IsaacSimVideoTrack", [
            "继承 VideoStreamTrack",
            "Replicator API 帧捕获",
            "Viewport Camera 降级方案",
            "分辨率自适应 + LANCZOS 缩放",
            "帧率控制 30fps + 暖机机制",
        ], PURPLE),
        ("CameraController", [
            "轨道相机 (orbit/pan/zoom)",
            "lookAt 矩阵手动构造",
            "每个实验独立摄像机预设",
            "延迟重调整 (对抗Kit重置)",
            "Viewport API + USD xform 双策略",
        ], GREEN),
        ("遥测循环 _telemetry_loop", [
            "100Hz 物理数据广播",
            "8个实验独立数据通道",
            "DynamicControl API 读速度",
            "碰撞检测 + 周期测量",
            "解析回退 (仿真数据兜底)",
        ], ORANGE),
    ]

    x = Inches(0.4)
    for title, items, color in backend_cards:
        add_card(slide, x, Inches(1.8), Inches(2.2), Inches(2.4), title, items, color)
        x += Inches(2.35)

    code_lines = [
        "class WebRTCServer:                          # 主服务器类",
        "  async def offer(self, request):             # WebRTC SDP 信令",
        "  async def _handle_ws_message(self, ws, d):  # 实验控制命令分发",
        "  async def _telemetry_loop(self):            # 100Hz 遥测广播",
        "  async def video_feed_handler(self, req):    # WS-JPEG 降级视频",
    ]
    y = Inches(4.5)
    for line in code_lines:
        add_text_box(slide, Inches(0.8), y, Inches(8.5), Inches(0.22),
                     line, font_size=9, color=BLUE_ACCENT, font_name="Consolas")
        y += Inches(0.24)

    add_page_number(slide, 11, TOTAL)

    # ========================================================
    # SLIDE 12: Section — 前端开发
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "04", "前端开发", "React 组件体系 · 实验UI · 学长代码适配")
    add_page_number(slide, 12, TOTAL)

    # ========================================================
    # SLIDE 13: 前端技术细节
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "前端组件架构", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    components = [
        ("Landing.tsx", "首页动画着陆页", RGBColor(0xEF, 0x44, 0x44)),
        ("LevelSelect.tsx", "8实验网格选择器 (锁定/解锁)", RGBColor(0xF5, 0x9E, 0x0B)),
        ("ExperimentView.tsx", "1237行! 实验控制+图表+视频+4试次", RGBColor(0x3B, 0x82, 0xF6)),
        ("WebRTCIsaacViewer.tsx", "WebRTC视频+WS-JPEG降级+相机交互", RGBColor(0x10, 0xB9, 0x81)),
        ("LabReportPDF.tsx", "实验一PDF报告React组件", RGBColor(0x8B, 0x5C, 0xF6)),
        ("Exp7ReportPDF.tsx", "实验七PDF报告React组件", RGBColor(0x14, 0xB8, 0xA6)),
    ]

    y = Inches(1.3)
    for name, desc, color in components:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), y, Inches(4.2), Inches(0.45))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"  {name}  —  {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Consolas"
        y += Inches(0.52)

    frontend_items = [
        ("experiments.ts — 声明式实验配置", [
            "每个实验: id, title, controls[], chartConfig[], extraMetrics[]",
            "controls: slider (min/max/step/command) / button 类型",
            "chartConfig: 遥测key映射 + 颜色 + Y轴ID",
            "isLocked: 未实现实验灰色锁定",
        ], CYAN),
        ("isaacService.ts — WebSocket 客户端", [
            "连接管理 + 自动重连",
            "命令发送: sendCommand(type, value)",
            "遥测回调订阅: onTelemetry()",
            "自定义消息分发 (exp2_progress等)",
        ], PURPLE),
        ("config.ts — 服务器URL自动检测", [
            "window.location.hostname 自动适配",
            "支持 .env 手动覆写",
            "HTTP :8080 / WS :30000 端口映射",
        ], GREEN),
    ]

    y = Inches(1.3)
    x_right = Inches(5.0)
    for title, items, color in frontend_items:
        add_text_box(slide, x_right, y, Inches(4.5), Inches(0.3),
                     "◆ " + title, font_size=11, color=color, bold=True)
        y += Inches(0.3)
        for item in items:
            add_text_box(slide, x_right + Inches(0.2), y, Inches(4.3), Inches(0.22),
                         "· " + item, font_size=9, color=GRAY_SUB)
            y += Inches(0.22)
        y += Inches(0.1)

    add_page_number(slide, 13, TOTAL)

    # ========================================================
    # SLIDE 14: 学长代码适配挑战
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "前端适配挑战 — 学长代码重构", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.5), Inches(0.5),
                 "前端基于学长 (Y-xvan/physical-lab) 项目进行了大量重构和适配工作",
                 font_size=13, color=GRAY_TEXT)

    challenges = [
        ("问题", "解决方案", RED_SOFT, GREEN),
        ("CDN Tailwind (index.html 直接引入)", "迁移到 PostCSS + Vite 插件构建", None, None),
        ("硬编码服务器IP (localhost:8080)", "config.ts 自动检测 hostname", None, None),
        ("只有 pause 命令 (服务器无 pause)", "映射 pause → stop_simulation", None, None),
        ("无实验切换机制", "添加 enter_experiment 协议 + switch_camera", None, None),
        ("单一实验UI (无法扩展)", "experiments.ts 声明式配置系统", None, None),
        ("无连接状态管理", "ConnectionStatus 枚举 + 自动重连", None, None),
        ("无数据导出功能", "CSV 导出 + PDF 报告 (react-pdf)", None, None),
        ("无实时遥测图表", "Recharts 实时折线图 + 300点滚动窗口", None, None),
        ("WebRTC 连接不稳定", "自动降级到 WS-JPEG (SSH隧道友好)", None, None),
        ("相机控制缺失", "鼠标拖拽: LMB=轨道 RMB=平移 滚轮=缩放", None, None),
    ]

    y = Inches(1.8)
    for i, (prob, sol, _, _) in enumerate(challenges):
        if i == 0:
            box_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(0.5), y, Inches(4.2), Inches(0.38))
            box_l.fill.solid()
            box_l.fill.fore_color.rgb = RED_SOFT
            box_l.line.fill.background()
            tf = box_l.text_frame
            p = tf.paragraphs[0]
            p.text = "  原始问题"
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(5.0), y, Inches(4.5), Inches(0.38))
            box_r.fill.solid()
            box_r.fill.fore_color.rgb = GREEN
            box_r.line.fill.background()
            tf = box_r.text_frame
            p = tf.paragraphs[0]
            p.text = "  我的解决方案"
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        else:
            bg_c = RGBColor(0xFE, 0xF2, 0xF2) if i % 2 == 1 else WHITE
            box_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(0.5), y, Inches(4.2), Inches(0.35))
            box_l.fill.solid()
            box_l.fill.fore_color.rgb = bg_c
            box_l.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            box_l.line.width = Pt(0.5)
            tf = box_l.text_frame
            p = tf.paragraphs[0]
            p.text = "  " + prob
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0xCC, 0x33, 0x33)
            p.font.name = "Microsoft YaHei"

            bg_c2 = RGBColor(0xF0, 0xFD, 0xF4) if i % 2 == 1 else WHITE
            box_r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(5.0), y, Inches(4.5), Inches(0.35))
            box_r.fill.solid()
            box_r.fill.fore_color.rgb = bg_c2
            box_r.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            box_r.line.width = Pt(0.5)
            tf = box_r.text_frame
            p = tf.paragraphs[0]
            p.text = "  ✓ " + sol
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0x16, 0x6B, 0x34)
            p.font.name = "Microsoft YaHei"

        y += Inches(0.38)

    add_page_number(slide, 14, TOTAL)

    # ========================================================
    # SLIDE 15: Section — 视频投送
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "05", "视频流投送", "Isaac Sim 渲染画面如何到达浏览器")
    add_page_number(slide, 15, TOTAL)

    # ========================================================
    # SLIDE 16: 视频管线详解
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "视频流管线 — 从GPU到浏览器", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    pipeline_steps = [
        ("1. GPU 渲染", "Isaac Sim\nPhysX 5 渲染", RGBColor(0x6B, 0x21, 0xA8)),
        ("2. 帧捕获", "Replicator API\nViewport Camera", RGBColor(0x26, 0x4E, 0x53)),
        ("3. 编码传输", "WebRTC Track\nH264 编码", RGBColor(0xE9, 0x72, 0x2D)),
        ("4. 浏览器解码", "<video> 元素\n实时播放", RGBColor(0x22, 0x7C, 0x9D)),
    ]

    x = Inches(0.3)
    for i, (step, desc, color) in enumerate(pipeline_steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, Inches(1.3), Inches(2.0), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(12)
        p.font.color.rgb = GOLD
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)

        if i < 3:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          x + Inches(2.05), Inches(1.7),
                                          Inches(0.35), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = GRAY_LIGHT
            arr.line.fill.background()
        x += Inches(2.4)

    details = [
        ("WebRTC 主通道 (优先)", [
            "SDP offer/answer 通过 HTTP POST /offer 交换",
            "ICE 候选地址补丁: 替换 0.0.0.0/127.0.0.1 为真实 LAN IP",
            "STUN 服务器: stun:stun.l.google.com:19302",
            "IsaacSimVideoTrack: 继承 aiortc.VideoStreamTrack",
            "30fps 帧率控制 + 30帧暖机 (等待渲染器就绪)",
            "Replicator render_product → rgb_annotator → numpy array → VideoFrame",
        ], CYAN),
        ("WS-JPEG 降级通道 (SSH隧道友好)", [
            "当 WebRTC ICE 连接失败 (8秒超时) 自动降级",
            "WebSocket /video_feed 端点接收 JPEG 帧",
            "_SharedFrameCapture 单例: 抓取视口 → PIL JPEG 编码 (quality=80)",
            "24fps 推送 + Blob URL 渲染到 <img> 元素",
            "前端自动检测: webrtc → ws-jpeg → none 三态切换",
        ], ORANGE),
        ("相机交互系统", [
            "鼠标拖拽: LMB=轨道旋转 RMB=平移 滚轮=缩放",
            "rAF 节流: requestAnimationFrame 合并相机指令",
            "fire-and-forget: fetch POST /camera (不等响应)",
            "方位角/仰角/距离 → lookAt矩阵 → USD Xform 应用",
        ], GREEN),
    ]

    y = Inches(2.8)
    for title, items, color in details:
        add_text_box(slide, Inches(0.5), y, Inches(4.5), Inches(0.3),
                     "◆ " + title, font_size=12, color=color, bold=True)
        y += Inches(0.3)
        for item in items:
            add_text_box(slide, Inches(0.9), y, Inches(8.5), Inches(0.22),
                         "· " + item, font_size=9, color=GRAY_SUB)
            y += Inches(0.22)
        y += Inches(0.1)

    add_page_number(slide, 16, TOTAL)

    # ========================================================
    # SLIDE 17: Section — 同学代码集成
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "06", "同学代码集成", "实验二（大摆角摆） · 实验七（动量守恒碰撞）")
    add_page_number(slide, 17, TOTAL)

    # ========================================================
    # SLIDE 18: 实验二集成
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "实验二 · 大摆角摆 — 同学代码集成", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    exp2_work = [
        ("同学原始代码", [
            "独立 Python 脚本 (expt2_large_amplitude_pendulum_sim_fixed.py)",
            "使用 isaacsim World 对象 + 程序化可视化摆杆",
            "RK4 积分求解阻尼复合摆微分方程",
            "零交叉法测量周期 + 级数展开比较",
            "matplotlib 绘图 + Markdown 报告输出",
        ], GRAY_SUB),
        ("我的集成工作", [
            "✦ 将 RK4 物理循环集成到 async 服务器 (_run_exp2_physics_loop)",
            "✦ 重写场景构建: World对象 → 纯 UsdGeom API (兼容 WebRTC 帧捕获)",
            "✦ 每6个RK4步骤 yield 一个Kit帧 (app.next_update_async) → ~360Hz 物理",
            "✦ 添加 θ/ω/α 遥测 + 实时周期检测 + 级数展开T理论值",
            "✦ 完整实验模式: 振幅扫描 + CSV + 5张图 + Markdown报告 → base64 推送到前端",
            "✦ 摄像机预设 + 延迟重调整 (1s/2s/4s 对抗 Kit 异步重置)",
            "✦ 前端: amplitude/damping 滑块 + 报告下载 (ZIP/MD/CSV)",
        ], CYAN),
    ]

    y = Inches(1.3)
    for title, items, color in exp2_work:
        add_text_box(slide, Inches(0.8), y, Inches(8.5), Inches(0.3),
                     "● " + title, font_size=13, color=color if color != GRAY_SUB else BLUE_DARK, bold=True)
        y += Inches(0.32)
        for item in items:
            bold = item.startswith("✦")
            add_text_box(slide, Inches(1.2), y, Inches(8), Inches(0.24),
                         item, font_size=10,
                         color=BLUE_MID if bold else GRAY_SUB)
            y += Inches(0.25)
        y += Inches(0.12)

    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(8.5), Inches(0.4),
                 "关键挑战: World 对象的 step() 与 WebRTC 帧捕获管线冲突 → 必须改用纯 USD + Kit async 渲染",
                 font_size=11, color=RED_SOFT, bold=True)

    add_page_number(slide, 18, TOTAL)

    # ========================================================
    # SLIDE 19: 实验七集成
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "实验七 · 动量守恒碰撞 — 同学代码集成", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    exp7_cards = [
        ("场景程序化构建", [
            "new_stage() → 纯 UsdGeom 构建",
            "零重力 (模拟水平无摩擦轨道)",
            "两个 DynamicCuboid 小车",
            "PhysicsMaterial: μ=0, e 可调",
            "PhysxRigidBodyAPI: CCD 碰撞",
            "64/32 求解器迭代 (高精度)",
        ], CYAN),
        ("碰撞物理管线", [
            "warmup 0.15s → 清零残余运动",
            "DynamicControl API 施加初速度",
            "轻量碰撞检测器 (速度突变法)",
            "碰撞后 0.4s 等待稳定",
            "超时兜底 (3s 无碰撞自动结束)",
            "碰撞前后速度精确捕获",
        ], PURPLE),
        ("前端4试次系统", [
            "与实验一相同的状态机模式",
            "idle → running → settled → recorded",
            "Cart1/Cart2 质量/速度滑块",
            "恢复系数 e 滑块 (0~1)",
            "实时 v1/v2/p_total 图表",
            "15列数据表 + CSV + PDF 报告",
        ], GREEN),
        ("遥测数据", [
            "v1, v2: 实时X方向速度",
            "p1, p2, p_total: 各车/总动量",
            "ke1, ke2, ke_total: 动能",
            "x1, x2: 小车位置",
            "phase: idle/warmup/running/settled",
            "碰撞类型: elastic/inelastic",
        ], ORANGE),
    ]

    x = Inches(0.3)
    for title, items, color in exp7_cards:
        add_card(slide, x, Inches(1.3), Inches(2.3), Inches(3.0), title, items, color)
        x += Inches(2.4)

    add_text_box(slide, Inches(0.8), Inches(4.6), Inches(8.5), Inches(0.4),
                 "关键挑战: PhysX 碰撞精度 → CCD + 高迭代求解器 + 零重力消除法向力干扰",
                 font_size=11, color=RED_SOFT, bold=True)

    add_page_number(slide, 19, TOTAL)

    # ========================================================
    # SLIDE 20: 技术挑战与解决
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "07", "技术挑战与解决", "Isaac Sim · WebRTC · 物理精度 · 工程问题")
    add_page_number(slide, 20, TOTAL)

    # ========================================================
    # SLIDE 21: 问题列表
    # ========================================================

    TOTAL = 22
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "踩过的坑与解决方案", font_size=28, color=BLUE_DARK, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.0), Inches(1.2), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = CYAN
    sep.line.fill.background()

    problems = [
        ("Isaac Sim 导入顺序", "SimulationApp 必须在 omni.*/pxr.* 之前实例化",
         "严格控制导入顺序; start_server.py 先 import isaacsim 再加载模块", CYAN),
        ("WebRTC ICE 连接失败", "SSH 隧道环境下 ICE 候选地址无法穿透",
         "SDP 补丁替换 IP + 8秒超时自动降级 WS-JPEG", PURPLE),
        ("摄像机被 Kit 重置", "切换实验后 Kit 异步初始化会覆盖摄像机",
         "延迟重调整策略: 1s/2s/4s 三次重设 + Viewport API + USD 双策略", GREEN),
        ("World.step() 与 WebRTC 冲突", "isaacsim World 的 step() 阻塞导致帧捕获失败",
         "放弃 World 对象 → 纯 UsdGeom + app.next_update_async() 异步渲染", ORANGE),
        ("物理暖机必要性", "world.reset() 后直接施加速度会被引擎吞掉",
         "0.15~1.0s 暖机阶段 (物理步进但保持初始位置)", RED_SOFT),
        ("VisualCuboid vs FixedCuboid", "装饰性几何体用 FixedCuboid 会产生幽灵碰撞",
         "严格区分: 装饰=VisualCuboid/UsdGeom.Cube, 碰撞=DynamicCuboid", TEAL),
        ("PhysX 碰撞精度不足", "小车碰撞穿透或速度丢失",
         "启用 CCD + solver iterations 64/32 + contactOffset 0.002", GOLD),
        ("Replicator 帧捕获不稳定", "首次使用时 rgb_annotator 返回空数据",
         "20帧预热 + KeyError 捕获 → 自动重建 render_product", BLUE_LIGHT),
    ]

    y = Inches(1.2)
    for title, problem, solution, color in problems:
        title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(0.4), y, Inches(2.0), Inches(0.55))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = color
        title_box.line.fill.background()
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)

        add_text_box(slide, Inches(2.5), y, Inches(3.3), Inches(0.55),
                     problem, font_size=9, color=RGBColor(0xCC, 0x33, 0x33))
        add_text_box(slide, Inches(6.0), y, Inches(3.7), Inches(0.55),
                     "→ " + solution, font_size=9, color=RGBColor(0x16, 0x6B, 0x34))
        y += Inches(0.62)

    add_page_number(slide, 21, TOTAL)

    # ========================================================
    # SLIDE 22: 成果与展望
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, BLUE_DARK, BLUE_MID)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5),
                 "项目成果与展望", font_size=32, color=WHITE, bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.1), Inches(1.5), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = GOLD
    sep.line.fill.background()

    achievements = [
        ("已完成", [
            "3 个实验完整实现（实验一全部自研 + 实验二/七集成适配）",
            "2100+ 行异步后端服务器（视频 + 遥测 + 8实验框架）",
            "1237 行前端组件（三种专用实验UI + 通用实验视图）",
            "双通道视频传输（WebRTC 优先 + WS-JPEG 降级）",
            "100Hz 实时遥测 + 实时数据可视化",
            "PDF 实验报告自动生成（前端 react-pdf + 后端 Jinja2）",
            "批处理 CLI 模式（CSV + matplotlib 图表 + Markdown 报告）",
            "完善的工程实践（配置中心化、文档体系、Agent连续性）",
        ], CYAN),
        ("未来展望", [
            "完成剩余 5 个实验（弹道摆/受驱阻尼/转动惯量/向心力/空气柱共振）",
            "VR 头显集成（reset() + 控制器事件 → core/vr.py）",
            "AI 智能实验编排（自动选参、自动运行、自动分析）",
            "多用户并发支持 + 教学管理后台",
        ], GOLD),
    ]

    y = Inches(1.5)
    for section_title, items, color in achievements:
        add_text_box(slide, Inches(0.8), y, Inches(8), Inches(0.4),
                     "◆ " + section_title, font_size=16, color=color, bold=True)
        y += Inches(0.4)
        for item in items:
            add_text_box(slide, Inches(1.2), y, Inches(8), Inches(0.28),
                         "✓ " + item, font_size=11, color=RGBColor(0xCC, 0xDD, 0xEE))
            y += Inches(0.28)
        y += Inches(0.2)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(8.5), Inches(0.5),
                 "感谢聆听  |  Q & A", font_size=24, color=GOLD,
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 22, TOTAL)

    # Fix total page numbers
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if "/20" in p.text:
                        p.text = p.text.replace("/20", f"/{len(prs.slides)}")

    out_path = os.path.join(os.path.dirname(__file__), "AI物理实验平台_项目汇报.pptx")
    prs.save(out_path)
    print(f"PPT saved to: {out_path}")
    return out_path


if __name__ == "__main__":
    build_ppt()
